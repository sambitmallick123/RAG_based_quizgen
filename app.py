import streamlit as st
import os
import json
import tempfile
from io import BytesIO
import numpy as np
from typing import List, Dict, Any
import uuid
import re

# Document processing libraries
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import pandas as pd

# ChromaDB for vector storage
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer

# Llama 3.2 model
import torch
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline
import warnings
warnings.filterwarnings('ignore')


def parse_number_ranges(input_str: str) -> List[int]:
    result = set()
    input_str = input_str.strip()
    if not input_str:
        return []
    parts = input_str.split(',')
    for part in parts:
        part = part.strip()
        if '-' in part:
            try:
                start, end = part.split('-')
                start = int(start)
                end = int(end)
                if start <= end:
                    result.update(range(start, end+1))
            except:
                pass
        else:
            try:
                num = int(part)
                result.add(num)
            except:
                pass
    return sorted(result)


def parse_sheet_names(input_str: str, available_sheets: List[str]) -> List[str]:
    result = set()
    input_str = input_str.strip()
    if not input_str:
        return []
    parts = input_str.split(',')
    for part in parts:
        part = part.strip()
        if '-' in part:
            try:
                start_s, end_s = part.split('-')
                start = int(start_s)
                end = int(end_s)
                if start <= end:
                    for i in range(start, end+1):
                        if 1 <= i <= len(available_sheets):
                            result.add(available_sheets[i-1])
            except:
                pass
        else:
            if part.isdigit():
                idx = int(part)
                if 1 <= idx <= len(available_sheets):
                    result.add(available_sheets[idx-1])
            else:
                matches = [sheet for sheet in available_sheets if sheet.lower() == part.lower()]
                if matches:
                    result.add(matches[0])
    return list(result)


class DocumentProcessor:
    @staticmethod
    def extract_from_pdf(file_content: BytesIO, pages: List[int] = None) -> str:
        try:
            reader = PyPDF2.PdfReader(file_content)
            text = ""
            total_pages = len(reader.pages)
            selected_pages = pages if pages else list(range(1, total_pages+1))
            selected_pages = [p for p in selected_pages if 1 <= p <= total_pages]
            for pagenum in selected_pages:
                page = reader.pages[pagenum-1]
                p_text = page.extract_text()
                if p_text:
                    text += p_text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PDF: {str(e)}")
            return ""

    @staticmethod
    def extract_from_docx(file_content: BytesIO, paragraphs: List[int] = None) -> str:
        try:
            doc = docx.Document(file_content)
            total_paragraphs = len(doc.paragraphs)
            selected_paragraphs = paragraphs if paragraphs else list(range(1, total_paragraphs + 1))
            selected_paragraphs = [p for p in selected_paragraphs if 1 <= p <= total_paragraphs]
            text = ""
            for pidx in selected_paragraphs:
                p_text = doc.paragraphs[pidx-1].text
                if p_text:
                    text += p_text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading DOCX: {str(e)}")
            return ""

    @staticmethod
    def extract_from_pptx(file_content: BytesIO, slides: List[int] = None) -> str:
        try:
            prs = Presentation(file_content)
            total_slides = len(prs.slides)
            selected_slides = slides if slides else list(range(1, total_slides+1))
            selected_slides = [s for s in selected_slides if 1 <= s <= total_slides]
            text = ""
            for slidenum in selected_slides:
                slide = prs.slides[slidenum-1]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PPTX: {str(e)}")
            return ""

    @staticmethod
    def extract_from_xlsx(file_content: BytesIO, sheets: List[str] = None) -> str:
        try:
            all_sheets_df = pd.read_excel(file_content, sheet_name=None)
            available_sheets = list(all_sheets_df.keys())
            selected_sheets = sheets if sheets else available_sheets
            selected_sheets = [s for s in selected_sheets if s in available_sheets]

            text = ""
            for sheet_name in selected_sheets:
                sheet_df = all_sheets_df[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for column in sheet_df.columns:
                    col_values = sheet_df[column].dropna().astype(str).values
                    if len(col_values) > 0:
                        text += f"{column}: {' '.join(col_values)}\n"
            return text
        except Exception as e:
            st.error(f"Error reading XLSX: {str(e)}")
            return ""


class ChromaVectorStore:
    def __init__(self, collection_name: str = "quiz_documents", persist_directory: str = "./chroma_db"):
        self.collection_name = collection_name
        self.persist_directory = persist_directory
        self.embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
        self.client = chromadb.PersistentClient(path=persist_directory)
        try:
            self.collection = self.client.get_collection(name=collection_name)
        except:
            self.collection = self.client.create_collection(
                name=collection_name,
                metadata={"hnsw:space": "cosine"}
            )
        self.chunks = []

    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        words = text.split()
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            if len(chunk.strip()) > 0:
                chunks.append(chunk.strip())
        return chunks

    def add_documents(self, texts: List[str], filenames: List[str] = None):
        all_chunks = []
        all_metadatas = []
        all_ids = []

        for idx, text in enumerate(texts):
            chunks = self.chunk_text(text)
            filename = filenames[idx] if filenames else f"document_{idx}"
            for chunk_idx, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                all_metadatas.append({
                    "source": filename,
                    "chunk_index": chunk_idx,
                    "total_chunks": len(chunks)
                })
                all_ids.append(f"{filename}_chunk_{chunk_idx}_{uuid.uuid4().hex[:8]}")

        if all_chunks:
            embeddings = self.embedding_model.encode(all_chunks).tolist()
            self.collection.add(
                documents=all_chunks,
                embeddings=embeddings,
                metadatas=all_metadatas,
                ids=all_ids
            )
            self.chunks = all_chunks
            st.success(f"Added {len(all_chunks)} chunks to ChromaDB collection")

    def retrieve_relevant_chunks(self, query: str, k: int = 3) -> List[str]:
        try:
            query_embedding = self.embedding_model.encode([query]).tolist()[0]
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=k
            )
            return results['documents'][0] if results['documents'] else []
        except Exception as e:
            st.error(f"Error retrieving chunks: {str(e)}")
            return []

    def get_collection_stats(self) -> Dict[str, Any]:
        try:
            count = self.collection.count()
            return {
                "total_documents": count,
                "collection_name": self.collection_name
            }
        except Exception as e:
            return {"error": str(e)}

    def clear_collection(self):
        try:
            self.client.delete_collection(name=self.collection_name)
            self.collection = self.client.create_collection(
                name=self.collection_name,
                metadata={"hnsw:space": "cosine"}
            )
            self.chunks = []
            st.success("Collection cleared successfully")
        except Exception as e:
            st.error(f"Error clearing collection: {str(e)}")


class QuizGenerator:
    def __init__(self):
        self.model_name = "/home/sambit/models/Llama-3.2-3B-Instruct"
        self.tokenizer = None
        self.model = None
        self.pipeline = None
        self.model_loaded = False
        self.load_model()

    def load_model(self):
        try:
            st.info("Loading Llama 3.2 model... This may take a few minutes.")
            self.tokenizer = AutoTokenizer.from_pretrained(self.model_name)
            if self.tokenizer.pad_token is None:
                self.tokenizer.pad_token = self.tokenizer.eos_token
            self.model = AutoModelForCausalLM.from_pretrained(
                self.model_name,
                torch_dtype=torch.bfloat16,
                device_map="auto",
                trust_remote_code=True,
                low_cpu_mem_usage=True
            )
            self.pipeline = pipeline(
                "text-generation",
                model=self.model,
                tokenizer=self.tokenizer,
                device_map="auto",
                torch_dtype=torch.bfloat16
            )
            self.model_loaded = True
            st.success("✅ Llama 3.2 model loaded successfully!")
        except Exception as e:
            st.error(f"Error loading model: {str(e)}")
            st.info("Please ensure you have access to meta-llama/Llama-3.2-3B-Instruct model and sufficient GPU memory")
            self.model_loaded = False

    def generate_mcq(self, context: str, question_num: int) -> Dict[str, Any]:
        if not self.model_loaded or not self.pipeline:
            return self._fallback_question(question_num)

        prompt = f"""Based on the following context, create a multiple-choice question with exactly 4 options (A, B, C, D) where only one option is correct.

Context: {context[:1000]}

Requirements:
- Create a clear, specific question based on the context
- Provide 4 distinct options labeled A, B, C, D
- Only one option should be correct
- Do not include explanations

Format your response exactly as:
Question: [Your question here]
A) [Option A]
B) [Option B] 
C) [Option C]
D) [Option D]
Correct Answer: [Option A, Option B, Option C, or Option D]

Generate the question now:"""

        try:
            messages = [
                {"role": "system", "content": "You are an expert quiz creator. Create clear, educational multiple-choice questions based on the provided context. Follow the exact format specified. Do not include explanations."},
                {"role": "user", "content": prompt}
            ]

            formatted_prompt = self.tokenizer.apply_chat_template(
                messages,
                tokenize=False,
                add_generation_prompt=True
            )

            response = self.pipeline(
                formatted_prompt,
                max_new_tokens=300,
                do_sample=True,
                temperature=0.7,
                top_p=0.9,
                pad_token_id=self.tokenizer.eos_token_id,
                eos_token_id=self.tokenizer.eos_token_id
            )

            generated_text = response[0]['generated_text']
            if formatted_prompt in generated_text:
                generated_text = generated_text.replace(formatted_prompt, "").strip()

            return self._parse_mcq_response(generated_text, question_num)

        except Exception as e:
            st.error(f"Error generating question {question_num}: {str(e)}")
            return self._fallback_question(question_num)

    def _parse_mcq_response(self, response: str, question_num: int) -> Dict[str, Any]:
        try:
            lines = [line.strip() for line in response.split('\n') if line.strip()]
            question = ""
            options = ["", "", "", ""]
            correct_answer_letter = "A"

            for line in lines:
                line_upper = line.upper()
                if line_upper.startswith("QUESTION:"):
                    question = line.split(":", 1)[1].strip()
                elif re.match(r"^[A-Da-d]\)", line):
                    idx = ord(line[0].upper()) - ord('A')
                    if 0 <= idx < 4:
                        options[idx] = line.split(")", 1)[1].strip()
                elif "CORRECT ANSWER" in line_upper:
                    after_colon = line.split(":", 1)[1] if ":" in line else line
                    match = re.search(r'\b([A-Da-d])\b', after_colon)
                    if match:
                        correct_answer_letter = match.group(1).upper()

            if not question:
                question = "Based on the provided context, which of the following is correct?"
            for i in range(4):
                if not options[i]:
                    options[i] = f"Option {chr(65+i)} (generated)"

            correct_answer_index = ord(correct_answer_letter) - ord('A')
            if not (0 <= correct_answer_index < 4):
                correct_answer_index = 0

            return {
                "question": question,
                "options": options,
                "correct_answer": correct_answer_index,
                "correct_option_letter": correct_answer_letter
            }

        except Exception as e:
            st.error(f"Error parsing response for question {question_num}: {str(e)}")
            return self._fallback_question(question_num)

    def _fallback_question(self, question_num: int) -> Dict[str, Any]:
        return {
            "question": f"What is the main topic discussed in the uploaded documents?",
            "options": [
                "Primary topic from the document",
                "Secondary topic mentioned",
                "Related concept discussed",
                "Supporting detail provided"
            ],
            "correct_answer": 0,
            "correct_option_letter": "A"
        }


class PlainTextFormatter:
    @staticmethod
    def format_quiz_to_plaintext(questions: List[Dict[str, Any]], quiz_title: str = "Generated MCQ Quiz") -> str:
        output = []
        output.append(f"Quiz: {quiz_title}")
        output.append("=" * 50)
        output.append("")
        for i, q in enumerate(questions, 1):
            output.append(f"{i}. {q['question']}")
            output.append("")
            for j, option in enumerate(q['options']):
                output.append(f"{chr(65+j)}. {option}")
            output.append("")
            output.append(f"Correct Answer: {q['correct_option_letter']}")
            output.append("")
            output.append("-" * 30)
            output.append("")
        return "\n".join(output)


def main():
    st.set_page_config(
        page_title="Quiz Generator",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    st.title("📚 RAG-Based Quiz Generator")
    st.markdown("Upload documents and generate multiple-choice questions using Llama 3.2 instruct model with ChromaDB vector storage")

    if 'vector_store' not in st.session_state:
        st.session_state.vector_store = ChromaVectorStore()
    if 'quiz_generator' not in st.session_state:
        with st.spinner("Initializing AI model..."):
            st.session_state.quiz_generator = QuizGenerator()
    if 'processor' not in st.session_state:
        st.session_state.processor = DocumentProcessor()

    with st.sidebar:
        st.header("📊 ChromaDB Collection")
        stats = st.session_state.vector_store.get_collection_stats()
        if "error" not in stats:
            st.metric("Total Documents", stats.get("total_documents", 0))
            st.text(f"Collection: {stats.get('collection_name', 'N/A')}")
        if st.button("🗑️ Clear Collection", help="Remove all documents from ChromaDB"):
            st.session_state.vector_store.clear_collection()
            st.rerun()
        st.divider()
        st.header("🤖 Model Status")
        model_status = "✅ Ready" if st.session_state.quiz_generator.model_loaded else "❌ Not Loaded"
        st.text(f"Llama 3.2: {model_status}")

    col1, col2 = st.columns([1, 1])
    with col1:
        st.header("📄 Document Upload")
        uploaded_files = st.file_uploader(
            "Upload documents",
            type=['pdf', 'docx', 'pptx', 'xlsx'],
            accept_multiple_files=True,
            help="Supported formats: PDF, Word documents, PowerPoint presentations, Excel spreadsheets"
        )
        if 'file_page_inputs' not in st.session_state:
            st.session_state['file_page_inputs'] = {}

        if uploaded_files:
            st.markdown("### Specify pages/slides/paragraphs/sheets for each uploaded file (optional):")
            for uploaded_file in uploaded_files:
                ext = uploaded_file.name.split('.')[-1].lower()
                filename = uploaded_file.name
                input_key = f"pages_{filename}"
                col_label = "Pages (e.g. 1,3-5)" if ext == "pdf" else \
                            ("Slides (e.g. 1,3-5)" if ext == "pptx" else \
                            ("Paragraphs (e.g. 1-5,8)" if ext == "docx" else \
                            "Sheets (names or numbers, e.g. Sheet1,2,4-5)"))
                page_input = st.text_input(
                    label=f"{filename} - {col_label}",
                    key=input_key,
                    placeholder="Leave empty to use entire document"
                )
                st.session_state['file_page_inputs'][filename] = page_input

            if st.button("📥 Process Documents", type="primary"):
                all_texts = []
                filenames_to_store = []
                for uploaded_file in uploaded_files:
                    ext = uploaded_file.name.split('.')[-1].lower()
                    filename = uploaded_file.name
                    page_input = st.session_state['file_page_inputs'].get(filename, '').strip()

                    if ext == "pdf":
                        pages = parse_number_ranges(page_input)
                        text = st.session_state.processor.extract_from_pdf(uploaded_file, pages if pages else None)
                    elif ext == "pptx":
                        slides = parse_number_ranges(page_input)
                        text = st.session_state.processor.extract_from_pptx(uploaded_file, slides if slides else None)
                    elif ext == "docx":
                        paragraphs = parse_number_ranges(page_input)
                        text = st.session_state.processor.extract_from_docx(uploaded_file, paragraphs if paragraphs else None)
                    elif ext == "xlsx":
                        try:
                            excel_file = uploaded_file
                            all_sheets_df = pd.read_excel(excel_file, sheet_name=None)
                            available_sheets = list(all_sheets_df.keys())
                        except Exception as e:
                            st.error(f"Error reading XLSX for sheets: {str(e)}")
                            available_sheets = []
                        sheets = parse_sheet_names(page_input, available_sheets) if page_input else None
                        uploaded_file.seek(0)
                        text = st.session_state.processor.extract_from_xlsx(uploaded_file, sheets if sheets else None)
                    else:
                        st.warning(f"Unsupported file type: {ext}")
                        text = ""

                    if text.strip():
                        all_texts.append(text)
                        filenames_to_store.append(filename)
                        st.success(f"✅ Processed: {filename}")
                    else:
                        st.warning(f"No text extracted from {filename}, check your selection or document content.")

                if all_texts:
                    st.session_state.vector_store.add_documents(all_texts, filenames_to_store)
                    st.rerun()

    with col2:
        st.header("⚙️ Quiz Configuration")
        col2a, col2b = st.columns([2, 1])
        with col2a:
            num_questions = st.slider(
                "Number of questions",
                min_value=1,
                max_value=20,
                value=5,
                help="Select how many MCQ questions to generate"
            )
        with col2b:
            quiz_title = st.text_input(
                "Quiz Title",
                value="AI Generated Quiz",
                help="Title for your quiz"
            )

        stats = st.session_state.vector_store.get_collection_stats()
        has_documents = stats.get("total_documents", 0) > 0

        if has_documents and st.session_state.quiz_generator.model_loaded:
            if st.button("🎯 Generate Quiz", type="primary"):
                with st.spinner("Generating MCQ questions using AI..."):
                    questions = []
                    progress_bar = st.progress(0)
                    for i in range(num_questions):
                        sample_query = f"educational content question {i+1} knowledge test"
                        relevant_chunks = st.session_state.vector_store.retrieve_relevant_chunks(
                            sample_query, k=2
                        )
                        context = " ".join(relevant_chunks) if relevant_chunks else "General knowledge question"
                        question = st.session_state.quiz_generator.generate_mcq(context, i+1)
                        questions.append(question)
                        progress_bar.progress((i+1) / num_questions)

                    st.session_state.generated_questions = questions
                    st.session_state.quiz_title = quiz_title
                    st.success(f"✅ Generated {num_questions} questions successfully!")
                    st.rerun()

        else:
            if not has_documents:
                st.warning("⚠️ Please upload and process documents first")
            if not st.session_state.quiz_generator.model_loaded:
                st.warning("⚠️ AI model is not loaded. Please check the model status in the sidebar.")

    if 'generated_questions' in st.session_state:
        st.header("📝 Generated Quiz Questions")
        for i, q in enumerate(st.session_state.generated_questions, 1):
            with st.expander(f"Question {i}: {q['question'][:80]}{'...' if len(q['question']) > 80 else ''}"):
                st.markdown(f"**{i}. {q['question']}**")
                st.markdown("")
                for j, option in enumerate(q['options']):
                    prefix = "✅" if j == q['correct_answer'] else "❌"
                    st.markdown(f"{prefix} **{chr(65+j)}.** {option}")
                st.markdown(f"**Correct Answer: {q['correct_option_letter']}**")

        st.header("📋 Plain Text Output")
        plaintext_formatter = PlainTextFormatter()
        plaintext_output = plaintext_formatter.format_quiz_to_plaintext(
            st.session_state.generated_questions,
            st.session_state.get('quiz_title', 'Generated Quiz')
        )

        col3, col4 = st.columns([3, 1])
        with col3:
            st.markdown("**Copy the plain text below:**")
            st.text_area(
                "Quiz Output",
                plaintext_output,
                height=400,
                help="Select all text and copy"
            )
        with col4:
            st.download_button(
                label="📥 Download Quiz",
                data=plaintext_output,
                file_name=f"{st.session_state.get('quiz_title', 'quiz').lower().replace(' ', '_')}.txt",
                mime="text/plain",
                help="Download the plain text quiz file"
            )

            if st.button("🗑️ Clear Questions"):
                if 'generated_questions' in st.session_state:
                    del st.session_state.generated_questions
                if 'quiz_title' in st.session_state:
                    del st.session_state.quiz_title
                st.rerun()


if __name__ == "__main__":
    main()

